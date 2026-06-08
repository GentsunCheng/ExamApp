"""
应用内文件查看器 - 多引擎组合方案
支持格式:
  - 文本类: txt, md, csv, log, json, xml, yaml, yml, rtf  → QTextBrowser
  - PDF    : pdf                                           → PyMuPDF 渲染为图片
  - 文档   : docx                                          → python-docx 提取文本
  - 幻灯片 : pptx                                          → python-pptx 提取文本和图片
  - 表格   : xlsx                                          → openpyxl 转 HTML 表格
  - 图片   : png, jpg, jpeg, gif, bmp                      → QLabel + QPixmap
  - 音视频 : mp4, avi, mov, mp3, wav 等                    → QMediaPlayer + QVideoWidget
  - 其他   : 降级为系统默认程序打开
"""

import os
import shutil
import logging
import html

import fitz
import markdown as md_lib
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

from PySide6.QtCore import Qt, QUrl, QTimer
from PySide6.QtGui import QPixmap, QImage, QFont, QDesktopServices
from PySide6.QtWidgets import (
    QDialog, QVBoxLayout, QHBoxLayout, QLabel, QTextBrowser,
    QScrollArea, QPushButton, QTabWidget, QWidget, QSizePolicy,
    QFileDialog, QSlider,
)
from PySide6.QtMultimedia import QMediaPlayer, QAudioOutput
from PySide6.QtMultimediaWidgets import QVideoWidget

from theme_manager import theme_manager
from language import tr
from utils import show_info, show_warn

logger = logging.getLogger(__name__)

# 支持格式分类
_TEXT_EXTS = {'.txt', '.csv', '.log', '.json', '.xml', '.yaml', '.yml', '.rtf', '.ini', '.cfg', '.py', '.sh', '.bat', '.html', '.htm'}
_MD_EXTS = {'.md'}
_IMAGE_EXTS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'}
_PDF_EXTS = {'.pdf'}
_DOCX_EXTS = {'.docx'}
_PPTX_EXTS = {'.pptx'}
_XLSX_EXTS = {'.xlsx'}
_MEDIA_EXTS = {'.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm', '.mp3', '.wav', '.ogg', '.flac', '.aac', '.wma', '.m4a'}
# 降级为系统打开的格式
_FALLBACK_EXTS = {'.doc', '.xls', '.ppt', '.zip', '.rar', '.7z'}


def is_supported_format(file_path):
    """判断文件是否可在应用内查看"""
    ext = os.path.splitext(file_path)[1].lower()
    return ext in _TEXT_EXTS | _MD_EXTS | _IMAGE_EXTS | _PDF_EXTS | _DOCX_EXTS | _PPTX_EXTS | _XLSX_EXTS | _MEDIA_EXTS


def open_file_in_viewer(file_path, parent=None, original_name=None):
    """
    统一入口：在应用内查看文件。
    不支持的格式降级为系统默认程序打开。
    """
    if not file_path or not os.path.exists(file_path):
        return
    ext = os.path.splitext(file_path)[1].lower()
    if is_supported_format(file_path):
        dlg = FileViewerDialog(file_path, parent, original_name=original_name)
        dlg.exec()
    else:
        # 降级为系统打开
        QDesktopServices.openUrl(QUrl.fromLocalFile(file_path))


class FileViewerDialog(QDialog):
    """应用内文件查看对话框"""

    def __init__(self, file_path, parent=None, original_name=None):
        super().__init__(parent)
        self.file_path = file_path
        self.filename = os.path.basename(file_path)
        self.original_name = original_name or self.filename
        self.ext = os.path.splitext(file_path)[1].lower()
        self.colors = theme_manager.get_theme_colors()

        self.setWindowTitle(self.original_name)
        self.setMinimumSize(800, 600)
        self.resize(1000, 750)
        self.setStyleSheet(
            f"QDialog {{ background-color: {self.colors['card_background']}; }}"
        )

        self._setup_ui()

    def _setup_ui(self):
        layout = QVBoxLayout()
        layout.setContentsMargins(12, 12, 12, 12)
        layout.setSpacing(8)

        # 文件名标题
        header = QLabel(self.original_name)
        header.setAlignment(Qt.AlignmentFlag.AlignLeft | Qt.AlignmentFlag.AlignVCenter)
        header.setStyleSheet(
            f"color: {self.colors['text_primary']}; font-size: 14px; font-weight: bold; padding: 4px 0;"
        )
        layout.addWidget(header)

        # 内容区域
        try:
            content_widget = self._build_content()
        except Exception as e:
            logger.exception("文件查看器渲染失败: %s", e)
            error_label = QLabel(tr('viewer.render_error', error=str(e)))
            error_label.setStyleSheet(f"color: {self.colors['error']}; padding: 40px;")
            error_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
            content_widget = error_label
        layout.addWidget(content_widget, 1)

        # 底部按钮
        btn_row = QHBoxLayout()
        btn_row.addStretch()

        btn_save_as = QPushButton(tr('viewer.save_as'))
        btn_save_as.setStyleSheet(
            f"QPushButton {{ background-color: {self.colors['success']}; color: #fff; "
            f"padding: 6px 16px; font-size: 13px; border-radius: 6px; }}"
        )
        btn_save_as.clicked.connect(self._save_as)
        btn_row.addWidget(btn_save_as)

        btn_open_external = QPushButton(tr('viewer.open_external'))
        btn_open_external.setStyleSheet(
            f"QPushButton {{ background-color: {self.colors['info']}; color: #fff; "
            f"padding: 6px 16px; font-size: 13px; border-radius: 6px; }}"
        )
        btn_open_external.clicked.connect(self._open_external)
        btn_row.addWidget(btn_open_external)

        btn_close = QPushButton(tr('common.close'))
        btn_close.setStyleSheet(
            f"QPushButton {{ background-color: {self.colors['primary']}; color: #fff; "
            f"padding: 6px 16px; font-size: 13px; border-radius: 6px; }}"
        )
        btn_close.clicked.connect(self.accept)
        btn_row.addWidget(btn_close)

        layout.addLayout(btn_row)
        self.setLayout(layout)

    def _build_content(self):
        if self.ext in _TEXT_EXTS:
            return self._build_text_viewer()
        elif self.ext in _MD_EXTS:
            return self._build_md_viewer()
        elif self.ext in _IMAGE_EXTS:
            return self._build_image_viewer()
        elif self.ext in _PDF_EXTS:
            return self._build_pdf_viewer()
        elif self.ext in _DOCX_EXTS:
            return self._build_docx_viewer()
        elif self.ext in _PPTX_EXTS:
            return self._build_pptx_viewer()
        elif self.ext in _XLSX_EXTS:
            return self._build_xlsx_viewer()
        elif self.ext in _MEDIA_EXTS:
            return self._build_media_viewer()
        else:
            label = QLabel(tr('viewer.unsupported_format', ext=self.ext))
            label.setAlignment(Qt.AlignmentFlag.AlignCenter)
            return label

    # ------------------------------------------------------------------ #
    #  文本类文件渲染
    # ------------------------------------------------------------------ #
    def _build_text_viewer(self):
        browser = QTextBrowser()
        browser.setOpenExternalLinks(False)
        bg = self.colors['card_background']
        fg = self.colors['text_primary']
        browser.setStyleSheet(
            f"QTextBrowser {{ background-color: {bg}; color: {fg}; "
            f"border: 1px solid {self.colors['border']}; border-radius: 6px; "
            f"padding: 8px; font-family: 'Menlo', 'Consolas', 'Courier New', monospace; font-size: 13px; }}"
        )
        try:
            with open(self.file_path, 'r', encoding='utf-8', errors='replace') as f:
                text = f.read()
        except Exception:
            with open(self.file_path, 'rb') as f:
                text = f.read().decode('utf-8', errors='replace')

        if self.ext in ('.html', '.htm'):
            browser.setHtml(text)
        else:
            # 转义 HTML 特殊字符后以 <pre> 方式渲染，保留格式
            escaped = html.escape(text)
            html_content = (
                f"<html><body style='background-color:{bg}; color:{fg};'>"
                f"<pre style='white-space:pre-wrap; word-wrap:break-word; "
                f"font-family:Menlo,Consolas,monospace; font-size:13px; margin:0;'>"
                f"{escaped}</pre></body></html>"
            )
            browser.setHtml(html_content)
        return browser

    # ------------------------------------------------------------------ #
    #  Markdown 文件渲染
    # ------------------------------------------------------------------ #
    def _build_md_viewer(self):
        try:
            with open(self.file_path, 'r', encoding='utf-8', errors='replace') as f:
                md_text = f.read()
        except Exception:
            with open(self.file_path, 'rb') as f:
                md_text = f.read().decode('utf-8', errors='replace')

        html_body = md_lib.markdown(
            md_text,
            extensions=[
                'fenced_code',     # ```code``` 代码块
                'codehilite',      # 代码高亮
                'tables',          # 表格
                'toc',             # 目录
                'nl2br',           # 换行转 <br>
                'sane_lists',      # 列表与缩进
            ],
            extension_configs={
                'codehilite': {
                    'css_class': 'highlight',
                    'guess_lang': True,
                },
            }
        )

        bg = self.colors['card_background']
        fg = self.colors['text_primary']
        border = self.colors['border']
        primary = self.colors['primary']
        code_bg = self.colors.get('background', '#f5f7fa')

        css = f"""
        <style>
            body {{
                background-color: {bg};
                color: {fg};
                font-family: -apple-system, 'Segoe UI', 'PingFang SC', 'Microsoft YaHei', sans-serif;
                font-size: 14px;
                line-height: 1.7;
                padding: 16px 24px;
                margin: 0;
            }}
            h1 {{ font-size: 1.8em; margin: 24px 0 12px; color: {fg}; border-bottom: 1px solid {border}; padding-bottom: 8px; }}
            h2 {{ font-size: 1.5em; margin: 20px 0 10px; color: {fg}; border-bottom: 1px solid {border}; padding-bottom: 6px; }}
            h3 {{ font-size: 1.25em; margin: 16px 0 8px; color: {fg}; }}
            h4 {{ font-size: 1.1em; margin: 12px 0 6px; color: {fg}; }}
            p {{ margin: 8px 0; }}
            a {{ color: {primary}; text-decoration: none; }}
            a:hover {{ text-decoration: underline; }}
            ul, ol {{ margin: 8px 0; padding-left: 24px; }}
            li {{ margin: 4px 0; }}
            blockquote {{
                margin: 12px 0;
                padding: 8px 16px;
                border-left: 4px solid {primary};
                background-color: {code_bg};
                color: {fg};
            }}
            code {{
                font-family: 'Menlo', 'Consolas', 'Courier New', monospace;
                background-color: {code_bg};
                padding: 2px 6px;
                border-radius: 4px;
                font-size: 13px;
                color: {primary};
            }}
            pre {{
                background-color: {code_bg};
                padding: 12px 16px;
                border-radius: 8px;
                border: 1px solid {border};
                overflow-x: auto;
                font-size: 13px;
                line-height: 1.5;
            }}
            pre code {{
                background: none;
                padding: 0;
                border-radius: 0;
                color: {fg};
            }}
            table {{
                border-collapse: collapse;
                margin: 12px 0;
                width: 100%;
                font-size: 13px;
            }}
            th {{
                background-color: {primary};
                color: #ffffff;
                padding: 8px 12px;
                text-align: center;
                font-weight: bold;
                border: 1px solid {border};
            }}
            td {{
                padding: 6px 12px;
                border: 1px solid {border};
                text-align: center;
            }}
            tr:nth-child(even) td {{
                background-color: {code_bg};
            }}
            hr {{
                border: none;
                border-top: 1px solid {border};
                margin: 20px 0;
            }}
            img {{
                max-width: 100%;
                border-radius: 4px;
            }}
            .toc {{
                background-color: {code_bg};
                padding: 12px 16px;
                border-radius: 8px;
                border: 1px solid {border};
                margin: 12px 0;
            }}
            .toc ul {{ list-style: none; padding-left: 16px; }}
            .highlight {{ background-color: transparent; }}
        </style>
        """

        browser = QTextBrowser()
        browser.setOpenExternalLinks(True)
        browser.setStyleSheet(
            f"QTextBrowser {{ background-color: {bg}; color: {fg}; "
            f"border: 1px solid {border}; border-radius: 6px; }}"
        )
        browser.setHtml(
            f"<html><head>{css}</head><body>{html_body}</body></html>"
        )
        return browser

    # ------------------------------------------------------------------ #
    #  图片文件渲染
    # ------------------------------------------------------------------ #
    def _build_image_viewer(self):
        scroll = QScrollArea()
        scroll.setWidgetResizable(False)
        scroll.setAlignment(Qt.AlignmentFlag.AlignCenter)
        scroll.setStyleSheet(
            f"QScrollArea {{ border: 1px solid {self.colors['border']}; border-radius: 6px; "
            f"background-color: {self.colors['background']}; }}"
        )

        pixmap = QPixmap(self.file_path)
        if pixmap.isNull():
            label = QLabel(tr('viewer.image_load_failed'))
            label.setAlignment(Qt.AlignmentFlag.AlignCenter)
            label.setStyleSheet(f"color: {self.colors['error']}; padding: 40px;")
            scroll.setWidget(label)
            return scroll

        label = QLabel()
        label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        label.setPixmap(pixmap)
        label.setSizePolicy(QSizePolicy.Policy.Ignored, QSizePolicy.Policy.Ignored)
        scroll.setWidget(label)
        return scroll

    # ------------------------------------------------------------------ #
    #  PDF 文件渲染 (PyMuPDF)
    # ------------------------------------------------------------------ #
    def _build_pdf_viewer(self):
        doc = fitz.open(self.file_path)
        scroll = QScrollArea()
        scroll.setWidgetResizable(False)
        scroll.setAlignment(Qt.AlignmentFlag.AlignHCenter | Qt.AlignmentFlag.AlignTop)
        scroll.setStyleSheet(
            f"QScrollArea {{ border: 1px solid {self.colors['border']}; border-radius: 6px; "
            f"background-color: {self.colors['background']}; }}"
        )

        container = QWidget()
        container_layout = QVBoxLayout()
        container_layout.setContentsMargins(8, 8, 8, 8)
        container_layout.setSpacing(8)
        container_layout.setAlignment(Qt.AlignmentFlag.AlignHCenter | Qt.AlignmentFlag.AlignTop)

        # DPI 缩放：使用 150 DPI 作为默认渲染精度
        zoom = 150.0 / 72.0
        mat = fitz.Matrix(zoom, zoom)

        page_count = doc.page_count
        for page_num in range(page_count):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(matrix=mat)

            # 将 pixmap 数据转为 QImage
            img = QImage(pix.samples, pix.width, pix.height, pix.stride, QImage.Format.Format_RGB888)
            qpix = QPixmap.fromImage(img.copy())

            page_label = QLabel()
            page_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
            page_label.setPixmap(qpix)
            page_label.setStyleSheet(
                f"background-color: {self.colors['card_background']}; "
                f"border: 1px solid {self.colors['border']}; border-radius: 4px;"
            )
            container_layout.addWidget(page_label)

            # 页码标注
            page_info = QLabel(f"— {page_num + 1} / {page_count} —")
            page_info.setAlignment(Qt.AlignmentFlag.AlignCenter)
            page_info.setStyleSheet(
                f"color: {self.colors['text_tertiary']}; font-size: 11px; padding: 2px 0;"
            )
            container_layout.addWidget(page_info)

        container_layout.addStretch()
        container.setLayout(container_layout)
        scroll.setWidget(container)
        doc.close()
        return scroll

    # ------------------------------------------------------------------ #
    #  DOCX 文件渲染 (python-docx)
    # ------------------------------------------------------------------ #
    def _build_docx_viewer(self):
        doc = Document(self.file_path)
        browser = QTextBrowser()
        bg = self.colors['card_background']
        fg = self.colors['text_primary']
        browser.setStyleSheet(
            f"QTextBrowser {{ background-color: {bg}; color: {fg}; "
            f"border: 1px solid {self.colors['border']}; border-radius: 6px; "
            f"padding: 16px 24px; font-size: 14px; line-height: 1.6; }}"
        )

        html_parts = []
        for para in doc.paragraphs:
            text = para.text or ''
            if not text.strip():
                html_parts.append('<br/>')
                continue
            style_name = (para.style.name or '').lower()
            if 'heading 1' in style_name:
                html_parts.append(f'<h1 style="color:{fg};">{_escape(text)}</h1>')
            elif 'heading 2' in style_name:
                html_parts.append(f'<h2 style="color:{fg};">{_escape(text)}</h2>')
            elif 'heading 3' in style_name:
                html_parts.append(f'<h3 style="color:{fg};">{_escape(text)}</h3>')
            else:
                html_parts.append(f'<p style="color:{fg}; margin:4px 0;">{_escape(text)}</p>')

        # 处理表格
        for table in doc.tables:
            html_parts.append('<table border="1" cellpadding="4" cellspacing="0" style="border-collapse:collapse; margin:8px 0; width:100%;">')
            for row in table.rows:
                html_parts.append('<tr>')
                for cell in row.cells:
                    cell_text = _escape(cell.text or '')
                    html_parts.append(f'<td style="padding:4px 8px; border:1px solid {self.colors["border"]};">{cell_text}</td>')
                html_parts.append('</tr>')
            html_parts.append('</table>')

        html_content = (
            f"<html><body style='background-color:{bg}; color:{fg};'>"
            + ''.join(html_parts)
            + "</body></html>"
        )
        browser.setHtml(html_content)
        return browser

    # ------------------------------------------------------------------ #
    #  PPTX 文件渲染 (python-pptx)
    # ------------------------------------------------------------------ #
    def _build_pptx_viewer(self):
        prs = Presentation(self.file_path)
        bg = self.colors['card_background']
        fg = self.colors['text_primary']
        border = self.colors['border']
        primary = self.colors['primary']

        container = QWidget()
        container_layout = QVBoxLayout()
        container_layout.setContentsMargins(8, 8, 8, 8)
        container_layout.setSpacing(12)
        container_layout.setAlignment(Qt.AlignmentFlag.AlignTop)

        for slide_num, slide in enumerate(prs.slides, start=1):
            # 幻灯片容器
            slide_widget = QWidget()
            slide_layout = QVBoxLayout()
            slide_layout.setContentsMargins(16, 12, 16, 12)
            slide_layout.setSpacing(6)

            slide_layout.addWidget(QLabel(
                f"<b style='color:{fg}; font-size:15px;'>{tr('viewer.pptx_slide', num=slide_num)}</b>"
            ))

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        font_size = None
                        is_bold = False
                        if para.runs:
                            font = para.runs[0].font
                            # 估算字号：PPTX 默认 18pt
                            try:
                                font_size = font.size.pt if font.size else None
                            except Exception:
                                font_size = None
                            is_bold = font.bold or False
                        size_px = font_size or 14
                        weight = 'bold' if is_bold else 'normal'
                        text_color = fg
                        label = QLabel(text)
                        label.setWordWrap(True)
                        label.setStyleSheet(
                            f"color: {text_color}; font-size: {size_px}px; "
                            f"font-weight: {weight}; padding: 1px 0;"
                        )
                        slide_layout.addWidget(label)

                elif shape.shape_type == 13:  # Picture
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        pixmap = QPixmap()
                        pixmap.loadFromData(image_bytes)
                        if not pixmap.isNull():
                            # 限制最大宽度 600px
                            if pixmap.width() > 600:
                                pixmap = pixmap.scaledToWidth(
                                    600, Qt.TransformationMode.SmoothTransformation
                                )
                            img_label = QLabel()
                            img_label.setPixmap(pixmap)
                            img_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
                            slide_layout.addWidget(img_label)
                    except Exception:
                        pass

            if slide_layout.count() > 1:
                slide_widget.setLayout(slide_layout)
                slide_widget.setStyleSheet(
                    f"background-color: {bg}; "
                    f"border: 1px solid {border}; border-radius: 8px; "
                    f"margin: 4px 0;"
                )
                container_layout.addWidget(slide_widget)

        container_layout.addStretch()
        container.setLayout(container_layout)

        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll.setWidget(container)
        scroll.setStyleSheet(
            f"QScrollArea {{ border: 1px solid {border}; border-radius: 6px; "
            f"background-color: {self.colors['background']}; }}"
        )
        prs.close()
        return scroll

    # ------------------------------------------------------------------ #
    #  音视频文件播放 (QMediaPlayer)
    # ------------------------------------------------------------------ #
    def _build_media_viewer(self):
        is_video = self.ext in {'.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm'}

        media_widget = QWidget()
        media_layout = QVBoxLayout()
        media_layout.setContentsMargins(0, 0, 0, 0)
        media_layout.setSpacing(4)

        # 播放器
        self._media_player = QMediaPlayer()
        self._audio_output = QAudioOutput()
        self._media_player.setAudioOutput(self._audio_output)

        if is_video:
            self._video_widget = QVideoWidget()
            self._video_widget.setStyleSheet(
                f"background-color: #000; border: 1px solid {self.colors['border']}; "
                f"border-radius: 6px;"
            )
            self._media_player.setVideoOutput(self._video_widget)
            self._video_widget.setMinimumHeight(400)
            media_layout.addWidget(self._video_widget, 1)
        else:
            # 音频：显示装饰信息
            info_label = QLabel(tr('viewer.audio_playing', name=self.original_name))
            info_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
            info_label.setStyleSheet(
                f"color: {self.colors['text_primary']}; font-size: 16px; "
                f"padding: 60px 20px; "
                f"border: 1px solid {self.colors['border']}; border-radius: 6px; "
                f"background-color: {self.colors['background']};"
            )
            media_layout.addWidget(info_label, 1)

        # 控制栏
        controls = QWidget()
        controls.setStyleSheet(
            f"background-color: {self.colors['card_background']}; "
            f"border: 1px solid {self.colors['border']}; border-radius: 6px;"
        )
        ctrl_layout = QVBoxLayout()
        ctrl_layout.setContentsMargins(8, 6, 8, 6)
        ctrl_layout.setSpacing(4)

        # 进度条
        self._media_slider = QSlider(Qt.Orientation.Horizontal)
        self._media_slider.setRange(0, 100)
        self._media_slider.setStyleSheet(
            f"QSlider::groove:horizontal {{ height: 4px; background: {self.colors['border']}; "
            f"border-radius: 2px; }}"
            f"QSlider::handle:horizontal {{ background: {self.colors['primary']}; "
            f"width: 14px; height: 14px; margin: -5px 0; border-radius: 7px; }}"
            f"QSlider::sub-page:horizontal {{ background: {self.colors['primary']}; "
            f"border-radius: 2px; }}"
        )
        self._media_slider.sliderMoved.connect(self._media_seek)
        ctrl_layout.addWidget(self._media_slider)

        # 按钮行
        btn_row = QHBoxLayout()
        btn_row.setSpacing(8)

        btn_play = QPushButton(tr('viewer.media_play'))
        btn_play.setStyleSheet(
            f"QPushButton {{ background-color: {self.colors['primary']}; color: #fff; "
            f"padding: 6px 16px; font-size: 13px; border-radius: 6px; }}"
        )
        btn_play.clicked.connect(self._media_toggle_play)
        btn_row.addWidget(btn_play)

        self._media_play_btn = btn_play

        # 时间标签
        self._media_time_label = QLabel("00:00 / 00:00")
        self._media_time_label.setStyleSheet(
            f"color: {self.colors['text_secondary']}; font-size: 12px;"
        )
        btn_row.addWidget(self._media_time_label)

        # 音量
        volume_slider = QSlider(Qt.Orientation.Horizontal)
        volume_slider.setRange(0, 100)
        volume_slider.setValue(50)
        volume_slider.setFixedWidth(100)
        volume_slider.setStyleSheet(
            f"QSlider::groove:horizontal {{ height: 4px; background: {self.colors['border']}; "
            f"border-radius: 2px; }}"
            f"QSlider::handle:horizontal {{ background: {self.colors['primary']}; "
            f"width: 12px; height: 12px; margin: -4px 0; border-radius: 6px; }}"
            f"QSlider::sub-page:horizontal {{ background: {self.colors['primary']}; "
            f"border-radius: 2px; }}"
        )
        volume_slider.valueChanged.connect(self._media_set_volume)
        btn_row.addWidget(volume_slider)

        btn_row.addStretch()
        ctrl_layout.addLayout(btn_row)
        controls.setLayout(ctrl_layout)
        media_layout.addWidget(controls)

        media_widget.setLayout(media_layout)

        # 连接信号
        self._media_player.positionChanged.connect(self._media_position_changed)
        self._media_player.durationChanged.connect(self._media_duration_changed)
        self._media_player.playbackStateChanged.connect(self._media_state_changed)

        # 设置媒体源并播放
        self._media_player.setSource(QUrl.fromLocalFile(self.file_path))
        self._media_player.play()

        # 窗口关闭时停止播放
        self.finished.connect(self._stop_media)

        return media_widget

    # ------------------------------------------------------------------ #
    #  媒体播放器辅助方法
    # ------------------------------------------------------------------ #
    def _media_toggle_play(self):
        if self._media_player.playbackState() == QMediaPlayer.PlaybackState.PlayingState:
            self._media_player.pause()
        else:
            self._media_player.play()

    def _media_seek(self, position):
        duration = self._media_player.duration()
        if duration > 0:
            self._media_player.setPosition(int(duration * position / 100))

    def _media_set_volume(self, value):
        self._audio_output.setVolume(value / 100.0)

    def _media_position_changed(self, position):
        duration = self._media_player.duration()
        if duration > 0:
            self._media_slider.blockSignals(True)
            self._media_slider.setValue(int(position * 100 / duration))
            self._media_slider.blockSignals(False)
        self._media_time_label.setText(
            f"{self._fmt_time(position)} / {self._fmt_time(self._media_player.duration())}"
        )

    def _media_duration_changed(self, duration):
        self._media_time_label.setText(
            f"{self._fmt_time(self._media_player.position())} / {self._fmt_time(duration)}"
        )

    def _media_state_changed(self, state):
        if state == QMediaPlayer.PlaybackState.PlayingState:
            self._media_play_btn.setText(tr('viewer.media_pause'))
        else:
            self._media_play_btn.setText(tr('viewer.media_play'))

    def _stop_media(self):
        self._media_player.stop()

    @staticmethod
    def _fmt_time(ms):
        if ms < 0:
            return "00:00"
        s = int(ms / 1000)
        m, s = divmod(s, 60)
        return f"{m:02d}:{s:02d}"

    # ------------------------------------------------------------------ #
    #  XLSX 文件渲染 (openpyxl)
    # ------------------------------------------------------------------ #
    def _build_xlsx_viewer(self):
        wb = load_workbook(self.file_path, data_only=True)
        tabs = QTabWidget()
        tabs.setStyleSheet(
            f"QTabWidget::pane {{ border: 1px solid {self.colors['border']}; border-radius: 6px; }}"
            f"QTabBar::tab {{ padding: 6px 16px; font-size: 13px; }}"
            f"QTabBar::tab:selected {{ color: {self.colors['primary']}; border-bottom: 2px solid {self.colors['primary']}; }}"
        )

        bg = self.colors['card_background']
        fg = self.colors['text_primary']
        border = self.colors['border']
        header_bg = self.colors['primary']
        alt_row = self.colors.get('background', '#f5f7fa')

        for ws in wb.worksheets:
            browser = QTextBrowser()
            browser.setStyleSheet(
                f"QTextBrowser {{ background-color: {bg}; color: {fg}; "
                f"border: none; padding: 8px; font-size: 13px; }}"
            )

            rows_html = []
            max_col = ws.max_column or 1
            max_row = ws.max_row or 1

            # 限制渲染行数防止大文件卡死
            render_limit = 500
            actual_rows = min(max_row, render_limit)

            for row_idx, row in enumerate(ws.iter_rows(min_row=1, max_row=actual_rows, max_col=max_col, values_only=True), start=1):
                row_bg = header_bg if row_idx == 1 else (alt_row if row_idx % 2 == 0 else bg)
                row_fg = '#ffffff' if row_idx == 1 else fg
                row_style = f"background-color:{row_bg}; color:{row_fg};"
                rows_html.append(f'<tr style="{row_style}">')
                for cell_val in row:
                    val = _escape(str(cell_val)) if cell_val is not None else ''
                    tag = 'th' if row_idx == 1 else 'td'
                    weight = 'bold' if row_idx == 1 else 'normal'
                    rows_html.append(
                        f'<{tag} style="padding:4px 8px; border:1px solid {border}; '
                        f'font-weight:{weight}; text-align:center; white-space:nowrap;">{val}</{tag}>'
                    )
                rows_html.append('</tr>')

            truncated_note = ''
            if max_row > render_limit:
                truncated_note = (
                    f'<p style="color:{self.colors["warning"]}; text-align:center; margin-top:8px;">'
                    f'{tr("viewer.xlsx_truncated", total=max_row, shown=render_limit)}</p>'
                )

            html_content = (
                f"<html><body style='background-color:{bg};'>"
                f"<table style='border-collapse:collapse; width:100%; font-size:13px;'>"
                + ''.join(rows_html)
                + "</table>"
                + truncated_note
                + "</body></html>"
            )
            browser.setHtml(html_content)

            scroll = QScrollArea()
            scroll.setWidgetResizable(True)
            scroll.setWidget(browser)
            scroll.setStyleSheet(f"QScrollArea {{ border: none; }}")

            tabs.addTab(scroll, ws.title or f'Sheet{row_idx}')

        wb.close()
        return tabs

    # ------------------------------------------------------------------ #
    #  保存 / 外部打开
    # ------------------------------------------------------------------ #
    def _save_as(self):
        """另存为：将文件以原始文件名导出到用户选择的位置"""
        save_path, _ = QFileDialog.getSaveFileName(
            self, tr('viewer.save_as'), self.original_name
        )
        if not save_path:
            return
        try:
            shutil.copy2(self.file_path, save_path)
            show_info(self, tr('common.success'), tr('viewer.save_as_success', path=save_path))
        except Exception as e:
            show_warn(self, tr('common.error'), str(e))

    def _open_external(self):
        QDesktopServices.openUrl(QUrl.fromLocalFile(self.file_path))


def _escape(text):
    """HTML 转义"""
    return html.escape(str(text))
